"""Binary document converters for LLM-oriented Markdown extraction."""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree

from .converter import html_to_markdown
from .ocr import OcrCollector, OcrOptions

PDF_MIME_TYPES = {"application/pdf", "application/x-pdf"}
DOCX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
PPTX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
}
XLSX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
DOCUMENT_MIME_TYPES = PDF_MIME_TYPES | DOCX_MIME_TYPES | PPTX_MIME_TYPES | XLSX_MIME_TYPES


@dataclass(frozen=True)
class DocumentConversion:
    markdown: str
    title: str | None = None


def is_supported_document(*, extension: str | None, mime_type: str | None) -> bool:
    ext = (extension or "").lower()
    mime = (mime_type or "").lower()
    return ext in DOCUMENT_EXTENSIONS or mime in DOCUMENT_MIME_TYPES


def convert_document_bytes(
    data: bytes,
    *,
    source_name: str,
    base_url: str,
    extension: str | None,
    mime_type: str | None,
    ocr: OcrOptions | None = None,
) -> DocumentConversion:
    ext = (extension or "").lower()
    mime = (mime_type or "").lower()

    if ext == ".pdf" or mime in PDF_MIME_TYPES:
        return _convert_pdf(data, ocr=ocr)
    if ext == ".docx" or mime in DOCX_MIME_TYPES:
        return _convert_docx(data, base_url=base_url, source_name=source_name, ocr=ocr)
    if ext == ".pptx" or mime in PPTX_MIME_TYPES:
        return _convert_pptx(data, ocr=ocr)
    if ext == ".xlsx" or mime in XLSX_MIME_TYPES:
        return _convert_xlsx(data, ocr=ocr)

    raise ValueError(
        "Unsupported binary document input. Supported document formats: PDF, DOCX, PPTX, XLSX."
    )


def _stringify_cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, date):
        return value.isoformat()
    return str(value)


def _markdown_table(rows: Iterable[Iterable[Any]]) -> str:
    materialized = [[_stringify_cell(cell) for cell in row] for row in rows]
    materialized = [row for row in materialized if any(cell.strip() for cell in row)]
    if not materialized:
        return ""

    width = max(len(row) for row in materialized)
    normalized = [row + [""] * (width - len(row)) for row in materialized]
    header = normalized[0]
    body = normalized[1:]

    def clean(value: str, fallback: str = "") -> str:
        value = re.sub(r"\s*\n\s*", " <br> ", value)
        value = re.sub(r"\s+", " ", value).strip()
        value = value.replace("|", "\\|")
        return value or fallback

    header = [clean(cell, f"Column {idx + 1}") for idx, cell in enumerate(header)]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(clean(cell, " ") for cell in row) + " |")
    return "\n".join(lines)


def _image_to_png_bytes(image: Any) -> bytes:
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _convert_pdf(data: bytes, *, ocr: OcrOptions | None = None) -> DocumentConversion:
    try:
        import pdfplumber
        from pdfminer.high_level import extract_text
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("PDF conversion requires pdfplumber and pdfminer.six.") from exc

    markdown = extract_text(io.BytesIO(data)).strip()
    ocr_collector = OcrCollector(ocr)
    try:
        table_chunks: list[str] = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables() or []
                for table_index, table in enumerate(tables, start=1):
                    table_md = _markdown_table(table)
                    if table_md:
                        table_chunks.append(
                            f"### Page {page_number} Table {table_index}\n\n{table_md}"
                        )
                page.close()

        if table_chunks:
            tables_md = "\n\n".join(table_chunks)
            markdown = f"{markdown}\n\n## Extracted PDF Tables\n\n{tables_md}".strip()
    except Exception:
        pass

    if ocr_collector.enabled:
        ocr_chunks: list[str] = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                page_ocr = _ocr_pdf_page_images(page, page_number, ocr_collector)
                if page_ocr:
                    ocr_chunks.extend(page_ocr)
                page.close()
        if ocr_chunks:
            ocr_md = "\n\n".join(ocr_chunks)
            markdown = f"{markdown}\n\n## Extracted Image OCR\n\n{ocr_md}".strip()

    return DocumentConversion(markdown=markdown)


def _ocr_pdf_page_images(page: Any, page_number: int, ocr_collector: OcrCollector) -> list[str]:
    chunks: list[str] = []
    images = page.images or []
    if not images:
        text = (page.extract_text() or "").strip()
        if len(text) < 50:
            try:
                png = _image_to_png_bytes(page.to_image(resolution=200).original)
            except Exception:
                return chunks
            block = ocr_collector.run(png, label=f"page {page_number}")
            if block:
                chunks.append(block)
        return chunks

    for image_index, image in enumerate(images, start=1):
        if not ocr_collector.enabled:
            break
        width = float(image.get("width") or 0)
        height = float(image.get("height") or 0)
        if width < 20 or height < 20:
            continue
        bbox = (
            float(image["x0"]),
            float(image["top"]),
            float(image["x1"]),
            float(image["bottom"]),
        )
        try:
            cropped = page.crop(bbox)
            png = _image_to_png_bytes(cropped.to_image(resolution=200).original)
        except Exception:
            continue
        block = ocr_collector.run(png, label=f"page {page_number} image {image_index}")
        if block:
            chunks.append(block)
    return chunks


def _docx_title(data: bytes) -> str | None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            core_xml = archive.read("docProps/core.xml")
    except (KeyError, zipfile.BadZipFile):
        return None

    try:
        root = ElementTree.fromstring(core_xml)
    except ElementTree.ParseError:
        return None

    ns = {"dc": "http://purl.org/dc/elements/1.1/"}
    title = root.findtext("dc:title", default="", namespaces=ns).strip()
    return title or None


def _convert_docx(
    data: bytes,
    *,
    base_url: str,
    source_name: str,
    ocr: OcrOptions | None = None,
) -> DocumentConversion:
    try:
        import mammoth
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("DOCX conversion requires mammoth.") from exc

    result = mammoth.convert_to_html(io.BytesIO(data))
    markdown = html_to_markdown(result.value, base_url)
    image_ocr = _ocr_docx_images(data, ocr)
    if image_ocr:
        markdown = f"{markdown}\n\n## Extracted Image OCR\n\n{image_ocr}"
    return DocumentConversion(
        markdown=markdown,
        title=_docx_title(data) or Path(source_name).name,
    )


def _ocr_docx_images(data: bytes, ocr: OcrOptions | None) -> str:
    ocr_collector = OcrCollector(ocr)
    if not ocr_collector.enabled:
        return ""

    chunks: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            image_names = sorted(
                name for name in archive.namelist() if name.startswith("word/media/")
            )
            for index, name in enumerate(image_names, start=1):
                block = ocr_collector.run(
                    archive.read(name),
                    label=f"docx image {index}",
                )
                if block:
                    chunks.append(block)
    except zipfile.BadZipFile:
        return ""
    return "\n\n".join(chunks)


def _safe_text(value: Any) -> str:
    text = getattr(value, "text", "") or ""
    return re.sub(r"\s+", " ", text).strip()


def _shape_position(shape: Any) -> tuple[int, int]:
    top = getattr(shape, "top", 0) or 0
    left = getattr(shape, "left", 0) or 0
    return int(top), int(left)


def _pptx_picture_alt(shape: Any) -> str:
    try:
        alt = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
    except Exception:
        alt = ""
    name = getattr(shape, "name", "") or "image"
    return re.sub(r"\s+", " ", alt or name).strip()


def _pptx_chart_markdown(chart: Any) -> str:
    try:
        title = _safe_text(chart.chart_title.text_frame) if chart.has_title else "Chart"
        categories = [category.label for category in chart.plots[0].categories]
        series = list(chart.series)
        rows: list[list[Any]] = [["Category"] + [item.name for item in series]]
        for idx, category in enumerate(categories):
            rows.append([category] + [item.values[idx] for item in series])
        table = _markdown_table(rows)
        return f"### {title}\n\n{table}" if table else f"### {title}"
    except Exception:
        return "[Unsupported chart]"


def _pptx_shape_markdown(
    shape: Any,
    *,
    title_shape: Any | None,
    ocr_collector: OcrCollector,
) -> list[str]:
    chunks: list[str] = []

    if getattr(shape, "has_table", False):
        rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
        table = _markdown_table(rows)
        if table:
            chunks.append(table)

    if getattr(shape, "has_chart", False):
        chunks.append(_pptx_chart_markdown(shape.chart))

    if getattr(shape, "has_text_frame", False):
        text = _safe_text(shape.text_frame)
        if text:
            prefix = "# " if shape == title_shape else ""
            chunks.append(prefix + text)

    if getattr(shape, "image", None) is not None:
        chunks.append(f"![{_pptx_picture_alt(shape)}](image)")
        block = ocr_collector.run(shape.image.blob, label=getattr(shape, "name", "pptx image"))
        if block:
            chunks.append(block)

    if hasattr(shape, "shapes"):
        for child in sorted(shape.shapes, key=_shape_position):
            chunks.extend(
                _pptx_shape_markdown(
                    child,
                    title_shape=title_shape,
                    ocr_collector=ocr_collector,
                )
            )

    return chunks


def _convert_pptx(data: bytes, *, ocr: OcrOptions | None = None) -> DocumentConversion:
    try:
        import pptx
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("PPTX conversion requires python-pptx.") from exc

    presentation = pptx.Presentation(io.BytesIO(data))
    ocr_collector = OcrCollector(ocr)
    deck_title: str | None = None
    slide_chunks: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        chunks = [f"<!-- Slide number: {slide_number} -->"]

        for shape in sorted(slide.shapes, key=_shape_position):
            chunks.extend(
                _pptx_shape_markdown(
                    shape,
                    title_shape=title_shape,
                    ocr_collector=ocr_collector,
                )
            )

        if title_shape is not None:
            title_text = _safe_text(title_shape.text_frame)
            if title_text and deck_title is None:
                deck_title = title_text

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = _safe_text(slide.notes_slide.notes_text_frame)
            if notes:
                chunks.append(f"### Notes\n\n{notes}")

        slide_chunks.append("\n\n".join(chunk for chunk in chunks if chunk.strip()))

    return DocumentConversion(markdown="\n\n".join(slide_chunks), title=deck_title)


def _convert_xlsx(data: bytes, *, ocr: OcrOptions | None = None) -> DocumentConversion:
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("XLSX conversion requires openpyxl.") from exc

    workbook = openpyxl.load_workbook(
        io.BytesIO(data),
        read_only=not bool(ocr and ocr.enabled),
        data_only=True,
    )
    ocr_collector = OcrCollector(ocr)
    chunks: list[str] = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        table = _markdown_table(rows)
        if table:
            chunks.append(f"## {sheet.title}\n\n{table}")
        image_blocks = _ocr_xlsx_images(sheet, ocr_collector)
        if image_blocks:
            chunks.append(f"### Images in {sheet.title}\n\n{image_blocks}")
    workbook.close()
    return DocumentConversion(markdown="\n\n".join(chunks))


def _ocr_xlsx_images(sheet: Any, ocr_collector: OcrCollector) -> str:
    if not ocr_collector.enabled:
        return ""

    chunks: list[str] = []
    for index, image in enumerate(getattr(sheet, "_images", []), start=1):
        if not ocr_collector.enabled:
            break
        try:
            data = image._data()
        except Exception:
            continue
        block = ocr_collector.run(data, label=f"{sheet.title} image {index}")
        if block:
            chunks.append(block)
    return "\n\n".join(chunks)

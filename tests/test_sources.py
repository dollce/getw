from __future__ import annotations

import tempfile
import unittest
import zipfile
from io import BytesIO
from pathlib import Path
from unittest.mock import patch

from click.testing import CliRunner

from mark2down.cleaner import clean_markdown
from mark2down.cli import main
from mark2down.converter import html_to_markdown
from mark2down.ocr import OcrOptions, format_ocr_block
from mark2down.sources import convert_bytes, convert_data_uri, convert_file_uri, convert_local_file


def _minimal_pdf(text: str) -> bytes:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 24 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = BytesIO()
    output.write(b"%PDF-1.4\n")
    offsets = [0]
    for idx, obj in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{idx} 0 obj\n".encode("ascii"))
        output.write(obj)
        output.write(b"\nendobj\n")
    xref_start = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.write(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.write(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_start}\n%%EOF\n"
        ).encode("ascii")
    )
    return output.getvalue()


def _minimal_docx(text: str) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "docProps/core.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/">
  <dc:title>DOCX Fixture</dc:title>
</cp:coreProperties>""",
        )
        archive.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>{text}</w:t></w:r></w:p>
  </w:body>
</w:document>""",
        )
    return buffer.getvalue()


class SourceConversionTests(unittest.TestCase):
    def test_converts_local_html_file_with_main_content(self) -> None:
        html = b"""
        <html>
          <head><title>Ignored title</title></head>
          <body>
            <nav>Navigation</nav>
            <article>
              <h1>Report</h1>
              <p>Important body text.</p>
              <img alt="chart" src="data:image/png;base64,AAAA" />
            </article>
          </body>
        </html>
        """
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "report.html"
            path.write_bytes(html)

            result = convert_local_file(path)

        self.assertEqual(result.source_type, "file")
        self.assertEqual(result.title, "Ignored title")
        self.assertIn("# Report", result.markdown)
        self.assertIn("Important body text.", result.markdown)
        self.assertNotIn("Navigation", result.markdown)
        self.assertIn("![chart](embedded-image)", result.markdown)
        self.assertNotIn("base64", result.markdown)

    def test_strips_article_sidebar_footer_and_comments(self) -> None:
        html = """
        <html>
          <body>
            <article>
              <header><h1>Threat report</h1></header>
              <div class="c-article__main">
                <div class="js-article-body">
                  <p>Primary article body.</p>
                </div>
              </div>
              <div class="c-article__footer">
                <a href="/tag/noise">Tag noise</a>
              </div>
              <div id="comments">
                <p>Comment form noise.</p>
              </div>
              <div class="c-article__sidebar">
                <h3>GReAT webinars</h3>
                <p>Related item noise.</p>
              </div>
            </article>
          </body>
        </html>
        """

        markdown = html_to_markdown(html, "https://securelist.example/report")

        self.assertIn("# Threat report", markdown)
        self.assertIn("Primary article body.", markdown)
        self.assertNotIn("Tag noise", markdown)
        self.assertNotIn("Comment form noise", markdown)
        self.assertNotIn("GReAT webinars", markdown)
        self.assertNotIn("Related item noise", markdown)

    def test_strips_syntax_highlighter_layout_tables(self) -> None:
        html = """
        <html>
          <body>
            <article>
              <p>Generated comments:</p>
              <div class="urvanov-syntax-highlighter-syntax">
                <textarea readonly>
 ✅ Port is now listening (no accepting)
 ❌ Port is already in use
 🔍 regsvr32.exe detected as parent. Attempting to terminate...</textarea>
                <table class="crayon-table">
                  <tr>
                    <td>1<br/>2<br/>3</td>
                    <td>✅ Port is now listening ( no accepting ) ❌ Port is already in use 🔍 regsvr32 . exe detected as parent . Attempting to terminate . . .</td>
                  </tr>
                </table>
              </div>
            </article>
          </body>
        </html>
        """

        markdown = html_to_markdown(html, "https://securelist.example/report")

        self.assertIn("✅ Port is now listening (no accepting)", markdown)
        self.assertIn("🔍 regsvr32.exe detected as parent. Attempting to terminate", markdown)
        self.assertNotIn("| 1 2 3 |", markdown)
        self.assertNotIn("| --- | --- |", markdown)

    def test_ocr_block_does_not_turn_parenthesized_text_into_link(self) -> None:
        markdown = clean_markdown(format_ocr_block("(powershell Get-CimInstance)", "url image 3"))

        self.assertIn("**Image OCR: url image 3**", markdown)
        self.assertIn("```text\n(powershell Get-CimInstance)", markdown)
        self.assertIn("\n```", markdown)
        self.assertNotIn("Image OCR: url image 3](powershell", markdown)

    def test_converts_csv_stdin_to_markdown_table(self) -> None:
        result = convert_bytes(
            b"name,score\nalpha,10\nbeta,20\n",
            source_name="-",
            source_type="stdin",
            base_url="",
            extension="csv",
        )

        self.assertIn("| name | score |", result.markdown)
        self.assertIn("| alpha | 10 |", result.markdown)
        self.assertEqual(result.extension, ".csv")

    def test_autodetects_csv_without_extension_hint(self) -> None:
        result = convert_bytes(
            b"name,score\nalpha,10\nbeta,20\n",
            source_name="-",
            source_type="stdin",
            base_url="",
        )

        self.assertIn("| name | score |", result.markdown)
        self.assertIn("| alpha | 10 |", result.markdown)
        self.assertEqual(result.extension, ".csv")

    def test_preserves_tsv_delimiter_when_extension_hint_is_provided(self) -> None:
        result = convert_bytes(
            b"name\tscore\nalpha\t10\nbeta\t20\n",
            source_name="-",
            source_type="stdin",
            base_url="",
            extension="tsv",
        )

        self.assertIn("| name | score |", result.markdown)
        self.assertIn("| alpha | 10 |", result.markdown)
        self.assertEqual(result.extension, ".tsv")

    def test_autodetects_pdf_without_extension_hint(self) -> None:
        result = convert_bytes(
            _minimal_pdf("Hello PDF"),
            source_name="download",
            source_type="file",
            base_url="",
        )

        self.assertIn("Hello PDF", result.markdown)
        self.assertEqual(result.extension, ".pdf")

    def test_converts_json_to_fenced_pretty_markdown(self) -> None:
        result = convert_bytes(
            b'{"b": 2, "a": 1}',
            source_name="payload.json",
            source_type="file",
            base_url="",
            extension=".json",
        )

        self.assertIn("```json", result.markdown)
        self.assertIn('"a": 1', result.markdown)
        self.assertIn('"b": 2', result.markdown)

    def test_converts_file_uri_and_data_uri_inputs(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "note.md"
            path.write_text("# Note\n\nBody", encoding="utf-8")

            file_result = convert_file_uri(path.as_uri())

        data_result = convert_data_uri("data:text/csv;charset=utf-8,name%2Cscore%0Aa%2C1")

        self.assertEqual(file_result.source_type, "file")
        self.assertIn("# Note", file_result.markdown)
        self.assertEqual(data_result.source_type, "data")
        self.assertIn("| name | score |", data_result.markdown)

    def test_converts_pdf_docx_pptx_and_xlsx_documents(self) -> None:
        pdf_result = convert_bytes(
            _minimal_pdf("Hello PDF"),
            source_name="sample.pdf",
            source_type="file",
            base_url="",
            extension=".pdf",
        )
        docx_result = convert_bytes(
            _minimal_docx("Hello DOCX"),
            source_name="sample.docx",
            source_type="file",
            base_url="",
            extension=".docx",
        )

        import openpyxl
        import pptx
        from PIL import Image, ImageDraw
        from pptx.util import Inches

        deck = pptx.Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[0])
        slide.shapes.title.text = "Deck Title"
        slide.placeholders[1].text = "Slide body"
        image = Image.new("RGB", (420, 120), "white")
        ImageDraw.Draw(image).text((20, 45), "OCR IMAGE", fill="black")
        image_buffer = BytesIO()
        image.save(image_buffer, format="PNG")
        image_buffer.seek(0)
        slide.shapes.add_picture(image_buffer, Inches(1), Inches(3))
        deck_buffer = BytesIO()
        deck.save(deck_buffer)
        with patch("mark2down.ocr.extract_text_from_image", return_value="OCR IMAGE TEXT"):
            pptx_result = convert_bytes(
                deck_buffer.getvalue(),
                source_name="deck.pptx",
                source_type="file",
                base_url="",
                extension=".pptx",
                ocr=OcrOptions(enabled=True),
            )

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Data"
        sheet.append(["name", "score"])
        sheet.append(["alpha", 10])
        workbook_buffer = BytesIO()
        workbook.save(workbook_buffer)
        xlsx_result = convert_bytes(
            workbook_buffer.getvalue(),
            source_name="workbook.xlsx",
            source_type="file",
            base_url="",
            extension=".xlsx",
        )

        self.assertIn("Hello PDF", pdf_result.markdown)
        self.assertIn("Hello DOCX", docx_result.markdown)
        self.assertEqual(docx_result.title, "DOCX Fixture")
        self.assertIn("# Deck Title", pptx_result.markdown)
        self.assertIn("Slide body", pptx_result.markdown)
        self.assertIn("OCR IMAGE TEXT", pptx_result.markdown)
        self.assertIn("## Data", xlsx_result.markdown)
        self.assertIn("| alpha | 10 |", xlsx_result.markdown)

    def test_cli_help_only_exposes_output_option(self) -> None:
        result = CliRunner().invoke(main, ["--help"])

        self.assertEqual(result.exit_code, 0, result.output)
        self.assertIn("-o, --output", result.output)
        for removed in (
            "--ocr",
            "--wait",
            "--timeout",
            "--header",
            "--stdout",
            "--no-save",
            "--filename",
            "--extension",
            "--mime-type",
            "--charset",
            "--install-browsers",
        ):
            self.assertNotIn(removed, result.output)

    def test_cli_autodetects_stdin_csv_and_writes_requested_file(self) -> None:
        runner = CliRunner()
        with runner.isolated_filesystem():
            result = runner.invoke(
                main,
                ["-", "-o", "data.md"],
                input=b"name,score\nalpha,10\n",
            )

            self.assertEqual(result.exit_code, 0, result.output)
            markdown = Path("data.md").read_text(encoding="utf-8")

        self.assertIn("source_type: stdin", markdown)
        self.assertIn("| name | score |", markdown)
        self.assertIn("| alpha | 10 |", markdown)


if __name__ == "__main__":
    unittest.main()
